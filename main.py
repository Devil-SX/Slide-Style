import os
from io import BytesIO
from pathlib import Path

import numpy as np
import pptx
from PIL import Image
from pptx import Presentation
from copy import deepcopy

os.environ["OMP_NUM_THREADS"] = "1"
from sklearn.cluster import KMeans

ORANGE = [255, 230, 204]
BLUE = [218, 232, 252]
# BLUE = [0, 0, 255]

my_color = np.array([[255, 230, 204], [218, 232, 252]])

pptx_file = "test.pptx"
output_dir = Path("output")
output_dir.mkdir(exist_ok=True)
output_file = output_dir / pptx_file

presentation = Presentation(pptx_file)
new_presentation = deepcopy(presentation)


k = 2
kmeans = KMeans(n_clusters=k, n_init=10, max_iter=100)

for i, slide in enumerate(presentation.slides):
    print(f"Processing slide {i}")
    remove_bias = 0
    for s,shape in enumerate(slide.shapes):
        if isinstance(shape, pptx.shapes.picture.Picture):
            img = Image.open(BytesIO(shape.image.blob)).convert("RGB")
            img = np.array(img)

            height = img.shape[0]
            width = img.shape[1]
            pixels = np.array(img).reshape(height * width, 3)

            gray_image = np.dot(pixels[..., :3], [0.2989, 0.5870, 0.1140])
            mask = (gray_image > 50) & (gray_image < 230)
            masked_pixels = pixels[mask].copy()


            kmeans.fit(masked_pixels)
            cluster_centers = kmeans.cluster_centers_.astype(int)
            compressed_image = pixels.copy()
            compressed_image[mask] = my_color[kmeans.labels_]

            offset = masked_pixels.astype(float) - cluster_centers[
                kmeans.labels_
            ].astype(float)
            image_revised = compressed_image.astype(float)

            image_revised[mask] += np.clip(offset, -10, 10)
            
            image_revised = image_revised.reshape(height, width, 3)
            image_revised = np.clip(image_revised, 0, 255).astype(np.uint8)

            img_bytesio = BytesIO()
            Image.fromarray(image_revised).save(img_bytesio, format="JPEG")
            img_bytes = img_bytesio.getvalue()

            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            new_presentation.slides[i].shapes.element.remove(new_presentation.slides[i].shapes[s+remove_bias].element)
            print("Remove!")
            remove_bias -= 1
            new_presentation.slides[i].shapes.add_picture(img_bytesio, left, top, width, height)

new_presentation.save(output_file)
